"""파일 → 텍스트 블록 추출.

모든 파서는 `list[Block]`을 돌려준다. Block은 "출처를 짚을 수 있는 최소 단위"다.
페이지/섹션/시트 이름을 여기서 붙여두면, 나중에 답변에 "○○.pdf 3페이지, §4.2"
같은 인용을 달 수 있다. 이게 RAG의 신뢰도를 좌우한다.
"""
from __future__ import annotations

import io
import json
import re
import subprocess
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path

from . import tabular


@dataclass
class Block:
    text: str
    page: int | None = None      # PDF 페이지 / 슬라이드 번호 (1-based)
    section: str | None = None   # 감지된 절 제목 (예: "4.2 Discussion")
    label: str | None = None     # 시트 이름, 셀 범위 등 부가 위치 정보


class ParseError(RuntimeError):
    """파싱이 실패했다 — 고쳐야 할 문제일 수 있다."""


class SkipFile(Exception):
    """의도적으로 제외한다 — 문제가 아니다.

    ParseError 와 구분하는 이유: '실패'로 집계되면 원인을 찾게 되는데,
    용량 초과 같은 정책적 제외는 찾을 원인이 없다. 통계에서 섞이면
    "고쳐야 할 것이 몇 개인지"를 알 수 없게 된다.
    """


# ---------------------------------------------------------------- 섹션 감지

# 논문/보고서에서 흔한 절 제목 패턴.
# 예: "3 Methods", "4.2. Discussion", "Ⅱ. 실험", "제 3 장", "ABSTRACT"
_SECTION_PATTERNS = [
    re.compile(r"^\s*(\d+(?:\.\d+)*)\.?\s+(\S.{0,80})$"),
    re.compile(r"^\s*([IVXivx]+|[ⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩ]+)\s*[.)]\s*(\S.{0,80})$"),
    re.compile(r"^\s*제\s*(\d+)\s*[장절]\s*(\S.{0,80})$"),
]
_BARE_HEADINGS = {
    "abstract", "introduction", "related work", "background", "method",
    "methods", "methodology", "materials and methods", "experiments",
    "experimental setup", "results", "results and discussion", "discussion",
    "conclusion", "conclusions", "references", "acknowledgments",
    "acknowledgements", "appendix",
    "초록", "서론", "관련연구", "관련 연구", "이론적 배경", "연구방법", "연구 방법",
    "실험", "실험방법", "실험 방법", "결과", "결과 및 고찰", "고찰", "결론",
    "참고문헌", "감사의 글", "부록",
}


# 절 번호로 인정할 최대값. 이보다 크면 연도("2011 창업")나 수량으로 본다.
_MAX_SECTION_NO = 30


def _plausible_section_number(num: str) -> bool:
    """'4.2'는 절 번호, '2011'은 연도. 이를 구분한다.

    점으로 나뉜 계층 번호(4.2, 1.1.3)는 그 자체로 절 번호라는 신호가 강하므로 통과시키고,
    단일 정수는 작은 값만 인정한다. 회의록 본문의 "2011 창업"이 절 제목으로
    오인되면 인용이 엉뚱한 곳을 가리키므로 이 구분이 필요하다.
    """
    if "." in num:
        return all(p.isdigit() and int(p) <= 99 for p in num.split(".") if p)
    if num.isdigit():
        return int(num) <= _MAX_SECTION_NO
    return True  # 로마숫자 등


def detect_section(line: str) -> str | None:
    """한 줄이 절 제목처럼 보이면 정규화한 제목을 돌려준다."""
    stripped = line.strip()
    if not stripped or len(stripped) > 100:
        return None
    if stripped.lower().rstrip(".:") in _BARE_HEADINGS:
        return stripped.rstrip(".:")
    for pat in _SECTION_PATTERNS:
        m = pat.match(stripped)
        if m:
            num, title = m.group(1), m.group(2).strip()
            if not _plausible_section_number(num):
                return None
            # 본문 문장이 숫자로 시작한 경우를 걸러낸다 (제목은 마침표로 끝나지 않음)
            if title.endswith((".", ",", ";")) or len(title.split()) > 12:
                return None
            return f"{num} {title}"
    return None


# ---------------------------------------------------------------- PDF

def parse_pdf(path: Path) -> list[Block]:
    import fitz  # pymupdf

    blocks: list[Block] = []
    try:
        doc = fitz.open(path)
    except Exception as e:
        raise ParseError(f"PDF 열기 실패: {e}") from e

    current_section: str | None = None
    with doc:
        for pno, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if not text.strip():
                continue  # 스캔본 페이지 — OCR 단계에서 따로 처리
            lines = text.splitlines()
            buf: list[str] = []
            for line in lines:
                sec = detect_section(line)
                if sec:
                    if buf:
                        blocks.append(Block("\n".join(buf), page=pno, section=current_section))
                        buf = []
                    current_section = sec
                buf.append(line)
            if buf:
                blocks.append(Block("\n".join(buf), page=pno, section=current_section))
    return blocks


def pdf_text_ratio(path: Path, sample: int = 10) -> float:
    """앞 몇 페이지에서 텍스트 레이어가 있는 페이지 비율.

    0에 가까우면 스캔본 → OCR 대상. 인덱싱 전에 이걸로 걸러낸다.
    """
    import fitz

    with fitz.open(path) as doc:
        n = min(len(doc), sample)
        if n == 0:
            return 0.0
        with_text = sum(1 for i in range(n) if doc[i].get_text("text").strip())
        return with_text / n


# ---------------------------------------------------------------- MS Office

# 일부 오래된 도구가 만든 docx는 관계 타입·네임스페이스 URI로 OOXML 표준화(2006) 이전의
# 초안 주소(purl.oclc.org/ooxml/...)를 그대로 쓴다 — 실측: 학회 제출용 논문 2건
# (Accepted_paper.docx, conference-template-letter_ICRA작성본260403.docx) 모두
# xmlns:w="http://purl.oclc.org/ooxml/wordprocessingml/main" 이었다. python-docx는
# 표준 주소(schemas.openxmlformats.org, .../2006/...)만 인식해서 officeDocument
# 관계를 못 찾고 KeyError 를 던진다 — 내용 자체는 정상 OOXML이라 주소만 바꾸면 읽힌다.
_LEGACY_OOXML_NS = re.compile(rb"purl\.oclc\.org/ooxml/([A-Za-z]+)/")


_MAX_OOXML_PATCH_BYTES = 300 * 1024 * 1024   # 압축 해제 총량 상한 (zip 폭탄 방지)


def _patch_legacy_ooxml_ns(path: Path) -> io.BytesIO:
    with zipfile.ZipFile(path) as src:
        total = sum(i.file_size for i in src.infolist())
        if total > _MAX_OOXML_PATCH_BYTES:
            raise ParseError(
                f"레거시 네임스페이스 복구를 위해 메모리에 올리기엔 너무 큼 "
                f"({total / 1048576:.0f}MB > {_MAX_OOXML_PATCH_BYTES // 1048576}MB)"
            )
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as dst:
            for item in src.infolist():
                data = src.read(item.filename)
                if item.filename.endswith((".xml", ".rels")):
                    data = _LEGACY_OOXML_NS.sub(rb"schemas.openxmlformats.org/\1/2006/", data)
                dst.writestr(item, data)
    buf.seek(0)
    return buf


def _list_level(para) -> int | None:
    """이 문단이 글머리 기호 목록의 몇 번째 들여쓰기 단계인지 (0부터 시작).
    목록의 일부가 아니면 None.
    """
    from docx.oxml.ns import qn

    pPr = para._p.find(qn("w:pPr"))
    if pPr is None:
        return None
    numPr = pPr.find(qn("w:numPr"))
    if numPr is None:
        return None
    ilvl = numPr.find(qn("w:ilvl"))
    if ilvl is None:
        return 0   # ilvl 생략은 0단계를 뜻한다
    val = ilvl.get(qn("w:val"))
    return int(val) if val is not None else 0


def parse_docx(path: Path) -> list[Block]:
    from docx import Document

    try:
        doc = Document(str(path))
    except KeyError:
        doc = Document(_patch_legacy_ooxml_ns(path))
    blocks: list[Block] = []
    current_section: str | None = None
    buf: list[str] = []

    def flush() -> None:
        if buf:
            blocks.append(Block("\n".join(buf), section=current_section))
            buf.clear()

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "") if para.style else ""
        # Word의 Heading 스타일을 절 경계로 신뢰한다 (정규식보다 정확하다)
        if style.startswith("Heading") or style in ("Title", "제목"):
            flush()
            current_section = text
            continue
        lvl = _list_level(para)
        if lvl is not None and lvl <= 1:
            # 회의록처럼 "이름 → 그 아래 업데이트 목록"으로 된 문서에서, 얕은
            # 글머리 항목(사람 이름 등)을 절 경계로 본다. 실측: 랩미팅 회의록이
            # 이 구조였는데 안 나누면 한 사람 업데이트를 찾는 질의가 같은 청크에
            # 섞인 다른 사람 업데이트에 묻혀 검색이 안 됐다.
            flush()
            current_section = text
            continue
        if (sec := detect_section(text)):
            flush()
            current_section = sec
        buf.append(text)
    flush()

    # 표는 별도 블록으로 — 본문과 섞으면 청킹이 망가진다
    for ti, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append(Block("\n".join(rows), section=current_section, label=f"표 {ti}"))
    return blocks


def parse_pptx(path: Path) -> list[Block]:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[Block] = []
    for sno, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        title = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title is None and shape == slide.shapes.title:
                title = text
            parts.append(text)
        # 발표자 노트에 실제 설명이 들어있는 경우가 많다.
        # has_notes_slide 가 True 여도 notes_text_frame 은 None 일 수 있다
        # (노트 슬라이드는 있는데 텍스트 틀이 없는 경우) — 실제로 여기서 터졌다.
        if slide.has_notes_slide:
            frame = slide.notes_slide.notes_text_frame
            notes = frame.text.strip() if frame is not None and frame.text else ""
            if notes:
                parts.append(f"[발표자 노트] {notes}")
        if parts:
            blocks.append(Block("\n".join(parts), page=sno, section=title,
                                label=f"슬라이드 {sno}"))
    return blocks


def _repair_xlsx_via_libreoffice(path: Path, timeout: int = 300) -> Path:
    """일부 xlsx는 `xl/styles.xml`에 openpyxl이 모르는 속성(`count`)이 있는
    `<cellStyle>`을 갖고 있어 `load_workbook()`이 `TypeError:
    CellStyle.__init__() got an unexpected keyword argument 'count'`로 실패한다.
    LibreOffice로 한 번 다시 저장하면 styles.xml이 표준 형태로 재생성돼 열린다.

    같은 디렉터리·같은 확장자로 변환하면 soffice가 원본을 잠근 채 덮어쓰려다
    실패한다(실측) — 그래서 하위 `_repaired/` 폴더에 결과를 낸다.
    """
    outdir = path.parent / "_repaired"
    repaired = outdir / path.name
    if repaired.exists() and repaired.stat().st_mtime >= path.stat().st_mtime:
        return repaired
    outdir.mkdir(exist_ok=True)
    parent = path.parent.resolve()
    proc = subprocess.run(
        ["docker", "run", "--rm", "-v", f"{parent}:/data",
         "--entrypoint", "soffice", LIBREOFFICE_IMAGE,
         "--headless", "--convert-to", "xlsx", "--outdir", f"/data/{outdir.name}",
         f"/data/{path.name}"],
        capture_output=True, timeout=timeout,
    )
    if (proc.returncode != 0 or not repaired.exists()
            or repaired.stat().st_mtime < path.stat().st_mtime):
        # 종료 코드 0이어도 실제로 산출물을 못 갱신했을 수 있다 — mtime으로 재확인
        # (stale 산출물이 이미 있던 경우 조용히 그걸 반환하는 사고를 막는다).
        err = proc.stderr.decode("utf-8", "replace").strip()
        raise ParseError(f"LibreOffice xlsx 복구 실패: {err[-300:] if err else '(출력 없음)'}")
    return repaired


def parse_xlsx(path: Path, max_rows_per_sheet: int = 2000) -> list[Block]:
    from openpyxl import load_workbook

    # read_only + data_only: 큰 파일에서 메모리를 아끼고 수식 대신 값을 읽는다
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except TypeError:
        wb = load_workbook(str(_repair_xlsx_via_libreoffice(path)),
                           read_only=True, data_only=True)
    blocks: list[Block] = []
    try:
        for ws in wb.worksheets:
            rows: list[str] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows_per_sheet:
                    rows.append(f"... (이하 생략: {ws.title} 시트가 {max_rows_per_sheet}행 초과)")
                    break
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                blocks.append(Block("\n".join(rows), section=ws.title, label=f"시트 {ws.title}"))
    finally:
        wb.close()
    return blocks


# ---------------------------------------------------------------- 한글 (HWPX)

_OLE2_MAGIC = b"\xd0\xcf\x11\xe0"


def parse_hwpml(path: Path) -> list[Block]:
    """HWPML — 한글 문서를 XML로 담는 별개 교환 포맷. .hwp 확장자를 달고 오지만
    OLE2 바이너리(v5)도 v3 바이너리도 아니다.

    실측: 법제처 국가법령정보센터가 배포하는 법령 문서(농림축산식품부훈령 등)가
    이 형식이었다 — `<HWPML>` 루트, 문단은 `<P>`, 실제 텍스트는 그 안의 `<CHAR>`에
    있다. hwp5txt 는 OLE2만 읽으므로 "Not an OLE2 Compound Binary File"로 실패하지만
    내용 자체는 온전한 XML이라 태그만 걷어내면 읽힌다.
    """
    import html

    data = path.read_text(encoding="utf-8", errors="replace")
    blocks: list[Block] = []
    for para in re.findall(r"<P\b[^>]*>(.*?)</P>", data, re.DOTALL):
        chars = re.findall(r"<CHAR[^>]*>(.*?)</CHAR>", para, re.DOTALL)
        text = "".join(html.unescape(c) for c in chars).strip()
        if text:
            blocks.append(Block(text))
    if not blocks:
        raise ParseError("HWPML이지만 <CHAR> 텍스트를 찾지 못함")
    return blocks


def parse_hwp5(path: Path, timeout: int = 120) -> list[Block]:
    """구형 .hwp (바이너리 HWP v5) — pyhwp 의 hwp5txt 로 텍스트를 뽑는다.

    지원할 값어치가 있는 이유: [Workspace]/LabMeeting 에만 .hwp 가 162개 있고,
    학회 초록·논문심사 보고서·업무협의 문서처럼 검색 수요가 높은 자료다.

    파이썬 API 대신 CLI 를 쓰는 이유: pyhwp 의 내부 API 는 문서화가 얕고
    버전마다 바뀐다. hwp5txt 는 이 패키지가 공식적으로 지원하는 인터페이스다.
    (참고: pyhwp 0.1b15 는 six 의존성이 빠져 있어 별도 설치가 필요하다.)

    .hwp 확장자라고 전부 OLE2(v5)는 아니다 — HWPML(XML)과 v3(pre-OLE2 바이너리)도
    같은 확장자를 쓴다. hwp5txt 에 넘기기 전에 매직 바이트로 걸러 사유를 정확히
    남긴다 (그러지 않으면 셋 다 "Not an OLE2 Compound Binary File"로 뭉뚱그려져
    "HWPML인데 파서가 없다"와 "손상됐다"를 구분할 수 없다).
    """
    head = path.open("rb").read(512)
    if not head.startswith(_OLE2_MAGIC):
        # BOM·XML 선언 유무·선언 앞 공백 등 변형이 있어도 잡히도록 앞부분만 본다
        # (선언 없이 <HWPML>로 바로 시작하는 문서도 실측에서 나올 수 있다).
        probe = head.lstrip(b"\xef\xbb\xbf").lstrip()
        if probe.startswith(b"<?xml") or probe[:8].lower().startswith(b"<hwpml"):
            return parse_hwpml(path)
        raise ParseError(
            "HWP v5(OLE2 바이너리)가 아님 — 구형 v3 등 hwp5txt가 지원하지 않는 "
            ".hwp 변종으로 보임"
        )

    exe = Path(sys.executable).parent / "hwp5txt"
    if not exe.exists():
        raise ParseError("hwp5txt 없음 — pip install pyhwp six")

    proc = subprocess.run([str(exe), str(path)], capture_output=True, timeout=timeout)
    if proc.returncode != 0:
        err = proc.stderr.decode("utf-8", "replace").strip().splitlines()
        raise ParseError(f"hwp5txt 실패: {err[-1] if err else '(사유 없음)'}")

    text = proc.stdout.decode("utf-8", "replace")
    if not text.strip():
        raise ParseError("hwp5txt 가 빈 텍스트를 반환 (암호 걸린 문서일 수 있음)")

    blocks: list[Block] = []
    current_section: str | None = None
    for chunk in re.split(r"\n{2,}", text):
        chunk = chunk.strip()
        if not chunk:
            continue
        if (sec := detect_section(chunk.splitlines()[0])):
            current_section = sec
        blocks.append(Block(chunk, section=current_section))
    return blocks


def parse_hwpx(path: Path) -> list[Block]:
    """HWPX는 ZIP + XML이라 표준 라이브러리로 읽을 수 있다."""
    import xml.etree.ElementTree as ET

    blocks: list[Block] = []
    current_section: str | None = None
    try:
        with zipfile.ZipFile(path) as zf:
            sections = sorted(n for n in zf.namelist()
                              if re.match(r"Contents/section\d+\.xml$", n))
            if not sections:
                raise ParseError("HWPX 안에 section XML이 없음")
            for sec_name in sections:
                root = ET.fromstring(zf.read(sec_name))
                # 네임스페이스가 버전마다 달라서 태그 지역명으로 찾는다
                paras: list[str] = []
                for elem in root.iter():
                    if elem.tag.rsplit("}", 1)[-1] == "t" and elem.text:
                        paras.append(elem.text)
                text = "".join(paras)
                for chunk in re.split(r"\n{2,}", text):
                    chunk = chunk.strip()
                    if not chunk:
                        continue
                    if (s := detect_section(chunk.splitlines()[0])):
                        current_section = s
                    blocks.append(Block(chunk, section=current_section))
    except zipfile.BadZipFile as e:
        raise ParseError(f"HWPX가 아니거나 손상됨 (구형 .hwp일 수 있음): {e}") from e
    return blocks


# ---------------------------------------------------------------- 평문 / 노트북 / HTML

# 평문/CSV 를 통째로 메모리에 올리는 것을 막는 상한.
# 실험 데이터 CSV 가 수백 MB 인 경우가 있는데, 그런 파일은 문서 검색 대상이 아니면서
# 파싱만으로 메모리를 다 먹고 청크 수만 폭발시킨다 (검색 품질도 떨어뜨린다).
MAX_TEXT_BYTES = 20 * 1024 * 1024      # 20MB
MAX_TEXT_LINES = 50_000


def parse_text(path: Path) -> list[Block]:
    # 숫자 표·로그면 본문을 청킹하지 않고 요약 한 덩이만 만든다.
    # 확장자나 크기가 아니라 내용(숫자 토큰 비율)으로 판단한다 — labrag/tabular.py 참고.
    if (prof := tabular.profile(path)) is not None:
        return [Block(tabular.summarize(path, prof, path.name), label="데이터 파일 요약")]

    size = path.stat().st_size
    if size > MAX_TEXT_BYTES:
        raise SkipFile(
            f"평문 파일이 너무 큼 ({size / 1048576:.1f}MB > "
            f"{MAX_TEXT_BYTES // 1048576}MB) — 수치 데이터로 판별되지도 않음"
        )
    raw = path.read_text(encoding="utf-8", errors="replace")
    blocks: list[Block] = []
    current_section: str | None = None
    buf: list[str] = []

    def flush() -> None:
        if buf:
            blocks.append(Block("\n".join(buf).strip(), section=current_section))
            buf.clear()

    lines = raw.splitlines()
    truncated = len(lines) > MAX_TEXT_LINES
    for line in lines[:MAX_TEXT_LINES]:
        if line.startswith("#"):  # 마크다운 헤딩
            flush()
            current_section = line.lstrip("#").strip()
            continue
        if (sec := detect_section(line)):
            flush()
            current_section = sec
        buf.append(line)
    flush()
    if truncated:
        blocks.append(Block(
            f"(이하 생략: 원본이 {len(lines):,}행으로 {MAX_TEXT_LINES:,}행 상한을 초과)",
            section=current_section))
    return [b for b in blocks if b.text]


def parse_notebook(path: Path) -> list[Block]:
    nb = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    blocks: list[Block] = []
    current_section: str | None = None
    for i, cell in enumerate(nb.get("cells", []), start=1):
        src = cell.get("source", [])
        text = "".join(src) if isinstance(src, list) else str(src)
        text = text.strip()
        if not text:
            continue
        if cell.get("cell_type") == "markdown":
            first = text.splitlines()[0]
            if first.startswith("#"):
                current_section = first.lstrip("#").strip()
            blocks.append(Block(text, section=current_section, label=f"셀 {i} (설명)"))
        else:
            blocks.append(Block(text, section=current_section, label=f"셀 {i} (코드)"))
    return blocks


def parse_html(path: Path) -> list[Block]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "lxml")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    text = soup.get_text("\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return [Block(p.strip()) for p in text.split("\n\n") if p.strip()]


# ---------------------------------------------------------------- 레거시 오피스 (.doc/.ppt)

# 구형 바이너리 .doc/.ppt는 python-docx/python-pptx가 아예 못 연다 (OOXML이 아니라
# 옛 OLE2 바이너리 포맷). 호스트에 LibreOffice가 없으므로(연구실_RAG_설치가이드 참고)
# 컨테이너의 soffice로 1회성 변환한다 — README "남은 작업"에 있던 항목.
LIBREOFFICE_IMAGE = "linuxserver/libreoffice"


def convert_legacy_office(path: Path, target_ext: str, timeout: int = 300) -> Path:
    """.doc/.ppt 를 같은 디렉터리에 .docx/.pptx 로 변환한다 (원본은 그대로 둔다).

    raw_dir 은 "내려받은/변환한 원본 캐시"로 설계돼 있어(config.py 참고) 변환 결과를
    원본 옆에 남겨도 된다 — 원본보다 오래된 변환 결과만 재사용하지 않고 다시 만든다
    (mtime 비교). 단, `fetch()`가 애초에 구버전 원본을 다시 받아오지 않으면 이 비교도
    무의미하다 — `fetch()`는 로컬 캐시가 있으면 드라이브 mod_time 과 무관하게 그대로
    돌려주는 기존 동작이라(labrag/fetch.py), 원본 자체가 갱신됐는지는 이 함수가 보장할
    수 없다.
    """
    converted = path.with_suffix(f".{target_ext}")
    if converted.exists() and converted.stat().st_mtime >= path.stat().st_mtime:
        return converted
    parent = path.parent.resolve()   # docker -v 는 절대 경로가 필요하다
    proc = subprocess.run(
        ["docker", "run", "--rm", "-v", f"{parent}:/data",
         "--entrypoint", "soffice", LIBREOFFICE_IMAGE,
         "--headless", "--convert-to", target_ext, "--outdir", "/data",
         f"/data/{path.name}"],
        capture_output=True, timeout=timeout,
    )
    if (proc.returncode != 0 or not converted.exists()
            or converted.stat().st_mtime < path.stat().st_mtime):
        # 종료 코드 0이어도 실제로 산출물을 못 갱신했을 수 있다 — mtime으로 재확인
        # (stale 산출물이 이미 있던 경우 조용히 그걸 반환하는 사고를 막는다).
        err = proc.stderr.decode("utf-8", "replace").strip()
        raise ParseError(f"LibreOffice 변환 실패: {err[-300:] if err else '(출력 없음)'}")
    return converted


def parse_doc_legacy(path: Path) -> list[Block]:
    return parse_docx(convert_legacy_office(path, "docx"))


def parse_ppt_legacy(path: Path) -> list[Block]:
    return parse_pptx(convert_legacy_office(path, "pptx"))


# ---------------------------------------------------------------- 디스패치

# 카테고리 → 파서. filetypes.categorize()가 준 값을 그대로 쓴다.
PARSERS = {
    "pdf": parse_pdf,
    "gslides": parse_pdf,     # Slides는 PDF로 export
    "docx": parse_docx,
    "gdoc": parse_docx,       # Docs는 docx로 export
    "doc_legacy": parse_doc_legacy,
    "pptx": parse_pptx,
    "ppt_legacy": parse_ppt_legacy,
    "xlsx": parse_xlsx,
    "gsheet": parse_xlsx,     # Sheets는 xlsx로 export
    "hwp": parse_hwpx,        # .hwpx만 — .hwp는 아래에서 걸러진다
    "text": parse_text,
    "table": parse_text,
    "code": parse_text,
    "notebook": parse_notebook,
}


def parse(path: Path, category: str) -> list[Block]:
    """카테고리에 맞는 파서를 골라 실행한다."""
    if category == "hwp":
        # .hwp(바이너리 v5)와 .hwpx(ZIP+XML)는 완전히 다른 포맷이다
        return parse_hwp5(path) if path.suffix.lower() == ".hwp" else parse_hwpx(path)
    if path.suffix.lower() in (".html", ".htm"):
        return parse_html(path)
    fn = PARSERS.get(category)
    if fn is None:
        raise ParseError(f"'{category}' 카테고리용 파서가 없음")
    return fn(path)
